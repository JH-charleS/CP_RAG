from __future__ import annotations

import argparse
import hashlib
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Iterable

import numpy as np
import pdfplumber
from docx import Document
from pptx import Presentation
from pymilvus import Collection, CollectionSchema, DataType, FieldSchema, connections, utility
from sentence_transformers import SentenceTransformer

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from core.config import get_settings

try:
    import fitz  # PyMuPDF
except ImportError:  # pragma: no cover
    fitz = None

try:
    from paddleocr import PaddleOCR
except ImportError:  # pragma: no cover
    PaddleOCR = None


SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".doc", ".docx", ".ppt", ".pptx"}


def iter_documents(source_dir: Path) -> Iterable[Path]:
    for path in source_dir.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path


def read_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def read_docx(path: Path) -> str:
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip())


def read_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                parts.append(text)
    return "\n".join(parts)


def read_pdf(path: Path) -> str:
    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
    return "\n".join(parts)


def convert_with_soffice(src: Path, out_ext: str, out_dir: Path) -> Path | None:
    command = [
        "soffice",
        "--headless",
        "--convert-to",
        out_ext,
        "--outdir",
        str(out_dir),
        str(src),
    ]
    try:
        subprocess.run(command, check=True, capture_output=True)
    except Exception:
        return None
    candidate = out_dir / f"{src.stem}.{out_ext}"
    return candidate if candidate.exists() else None


def read_with_conversion(path: Path) -> str:
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        converted = convert_with_soffice(path, "txt", temp_path)
        if converted is None:
            return ""
        return read_txt(converted)


def split_text(text: str, chunk_size: int = 900, overlap: int = 150) -> list[str]:
    normalized = " ".join(text.split())
    if not normalized:
        return []
    chunks: list[str] = []
    start = 0
    length = len(normalized)
    while start < length:
        end = min(start + chunk_size, length)
        chunks.append(normalized[start:end])
        if end >= length:
            break
        start = max(0, end - overlap)
    return chunks


class OCRExtractor:
    def __init__(self) -> None:
        self._ocr = None

    @property
    def available(self) -> bool:
        return self._ocr is not None and fitz is not None

    def ensure_initialized(self) -> bool:
        if self._ocr is not None:
            return True
        if PaddleOCR is None or fitz is None:
            return False
        try:
            print("[OCR] initializing PaddleOCR ...")
            self._ocr = PaddleOCR(use_angle_cls=True, lang="ch")
            print("[OCR] PaddleOCR ready.")
            return True
        except Exception as exc:
            print(f"[OCR] initialization failed: {exc}")
            self._ocr = None
            return False

    def _ocr_image_array(self, image: np.ndarray) -> str:
        assert self._ocr is not None
        result = self._ocr.ocr(image, cls=True)
        if not result:
            return ""
        lines: list[str] = []
        for item in result:
            for block in item:
                if len(block) >= 2 and len(block[1]) >= 1:
                    lines.append(str(block[1][0]))
        return "\n".join(lines)

    def ocr_pdf(self, path: Path, max_pages: int = 20) -> str:
        if not self.ensure_initialized():
            return ""
        assert fitz is not None
        doc = fitz.open(str(path))
        parts: list[str] = []
        for idx, page in enumerate(doc):
            if idx >= max_pages:
                break
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            arr = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
            text = self._ocr_image_array(arr)
            if text.strip():
                parts.append(text)
        doc.close()
        return "\n".join(parts)


def extract_text(path: Path, ocr_extractor: OCRExtractor, min_length: int = 80) -> str:
    ext = path.suffix.lower()
    if ext == ".txt":
        text = read_txt(path)
    elif ext == ".docx":
        text = read_docx(path)
    elif ext == ".pptx":
        text = read_pptx(path)
    elif ext == ".pdf":
        text = read_pdf(path)
    elif ext in {".doc", ".ppt"}:
        text = read_with_conversion(path)
    else:
        text = ""

    if len(text.strip()) >= min_length:
        return text

    # 方案B：文本抽取不足时回退 OCR
    if ext == ".pdf":
        ocr_text = ocr_extractor.ocr_pdf(path)
        if len(ocr_text.strip()) > len(text.strip()):
            return ocr_text
    elif ext in {".doc", ".docx", ".ppt", ".pptx"}:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            converted_pdf = convert_with_soffice(path, "pdf", temp_path)
            if converted_pdf:
                ocr_text = ocr_extractor.ocr_pdf(converted_pdf)
                if len(ocr_text.strip()) > len(text.strip()):
                    return ocr_text
    return text


def ensure_collection(name: str, dim: int) -> Collection:
    if utility.has_collection(name):
        return Collection(name)

    fields = [
        FieldSchema(name="doc_id", dtype=DataType.VARCHAR, max_length=256, is_primary=True, auto_id=False),
        FieldSchema(name="title", dtype=DataType.VARCHAR, max_length=1024),
        FieldSchema(name="content", dtype=DataType.VARCHAR, max_length=65535),
        FieldSchema(name="source", dtype=DataType.VARCHAR, max_length=2048),
        FieldSchema(name="file_type", dtype=DataType.VARCHAR, max_length=16),
        FieldSchema(name="embedding", dtype=DataType.FLOAT_VECTOR, dim=dim),
    ]
    schema = CollectionSchema(fields=fields, description="CP_RAG raw document chunks")
    collection = Collection(name=name, schema=schema)
    collection.create_index(
        field_name="embedding",
        index_params={"index_type": "AUTOINDEX", "metric_type": "IP", "params": {}},
    )
    return collection


def build_doc_id(path: Path, chunk_idx: int) -> str:
    digest = hashlib.md5(str(path).encode("utf-8"), usedforsecurity=False).hexdigest()[:10]
    return f"{path.stem[:120]}_{digest}_{chunk_idx}"


def main() -> None:
    parser = argparse.ArgumentParser(description="Ingest raw docs into Milvus collection cp_rag_document.")
    parser.add_argument(
        "--source",
        type=Path,
        default=Path("source_data/raw_data"),
        help="Directory containing raw source documents.",
    )
    parser.add_argument("--collection", type=str, default="cp_rag_document", help="Milvus collection name.")
    parser.add_argument("--batch-size", type=int, default=64, help="Batch size for Milvus insert.")
    parser.add_argument("--max-files", type=int, default=0, help="Optional cap for processed files (0 means all).")
    args = parser.parse_args()

    settings = get_settings()
    connections.connect(
        alias="default",
        host=settings.milvus_host,
        port=settings.milvus_port,
        db_name=settings.milvus_db_name,
        timeout=settings.milvus_timeout,
        secure=settings.milvus_secure,
        user=settings.milvus_user or None,
        password=settings.milvus_password or None,
    )

    source_dir = args.source.resolve()
    if not source_dir.exists():
        raise FileNotFoundError(f"Source directory not found: {source_dir}")

    embedder = SentenceTransformer(settings.embedding_model_name)
    sample_vec = embedder.encode("dimension probe", normalize_embeddings=True)
    dim = int(sample_vec.shape[0])

    collection = ensure_collection(args.collection, dim)
    ocr_extractor = OCRExtractor()

    docs = list(iter_documents(source_dir))
    if args.max_files > 0:
        docs = docs[: args.max_files]
    if not docs:
        print(f"No supported documents found in {source_dir}")
        return
    print(f"Found {len(docs)} files in {source_dir}")

    rows: list[dict[str, object]] = []
    for index, doc_path in enumerate(docs, start=1):
        print(f"[{index}/{len(docs)}] parsing: {doc_path.name}")
        try:
            text = extract_text(doc_path, ocr_extractor=ocr_extractor)
        except Exception as exc:
            print(f"[skip] {doc_path} parse failed: {exc}")
            continue
        if not text.strip():
            continue
        chunks = split_text(text)
        if not chunks:
            continue
        embeddings = embedder.encode(chunks, normalize_embeddings=True)
        for idx, (chunk, emb) in enumerate(zip(chunks, embeddings)):
            rows.append(
                {
                    "doc_id": build_doc_id(doc_path, idx),
                    "title": doc_path.stem[:1000],
                    "content": chunk[:65000],
                    "source": str(doc_path.relative_to(source_dir)),
                    "file_type": doc_path.suffix.lower(),
                    "embedding": [float(x) for x in emb.tolist()],
                }
            )
        if len(rows) >= args.batch_size:
            collection.insert(rows)
            print(f"Inserted {len(rows)} rows ...")
            rows.clear()

    if rows:
        collection.insert(rows)
        print(f"Inserted {len(rows)} rows ...")
    collection.flush()
    collection.load()
    print(f"Collection: {args.collection}, num_entities={collection.num_entities}")


if __name__ == "__main__":
    main()
